"""
PowerPoint Recap Generator

Generates a two-slide recap presentation in PowerPoint matching the Arize style:
- Slide 1: Key Initiatives, Challenges, Solution Requirements
- Slide 2: Questions for Next Call (highlighted in magenta)
"""

import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from models import RecapSlideData


# Color scheme matching the reference slide
MAGENTA = RGBColor(200, 30, 120)  # Magenta/pink for titles and questions
DARK_TEXT = RGBColor(50, 50, 60)  # Dark text
LIGHT_BG = RGBColor(245, 248, 255)  # Light blue-ish background


def _add_footer(slide):
    """Add Arize branding footer to a slide."""
    footer_box = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(3), Inches(0.3))
    tf = footer_box.text_frame
    tf.paragraphs[0].text = "Arize  |  We Make Models Work"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = MAGENTA


def _set_slide_background(slide):
    """Set the light blue-gray background for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def create_recap_presentation(data: RecapSlideData) -> bytes:
    """
    Create a PowerPoint presentation with two slides:
    - Slide 1: Recap (Key Initiatives, Challenges, Solution Requirements)
    - Slide 2: Questions for Next Call
    
    Args:
        data: RecapSlideData containing the content for the slides
        
    Returns:
        Bytes of the PowerPoint file
    """
    # Create presentation with widescreen dimensions
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    # =========================================
    # SLIDE 1: RECAP
    # =========================================
    slide1 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide1)
    
    # Title - Customer Name Recap
    title_text = f"{data.customer_name} Recap:" if data.customer_name else "Call Recap:"
    
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.paragraphs[0].text = title_text
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = MAGENTA
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle - Call Date (if available)
    if data.call_date:
        date_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.3))
        date_frame = date_box.text_frame
        date_frame.paragraphs[0].text = f"Call Date: {data.call_date}"
        date_frame.paragraphs[0].font.size = Pt(14)
        date_frame.paragraphs[0].font.italic = True
        date_frame.paragraphs[0].font.color.rgb = DARK_TEXT
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Section dimensions for 2-column layout
    col_width = Inches(6.0)
    left_x = Inches(0.5)
    right_x = Inches(6.833)
    top_y = Inches(1.1)
    
    # Helper function to add section title
    def add_section_title(slide, x, y, title, width=col_width):
        title_box = slide.shapes.add_textbox(x, y, width, Inches(0.4))
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = MAGENTA
    
    # Helper function to add bullet points
    def add_bullet_points(slide, x, y, items, width=col_width, highlight_all=False, max_items=5):
        content_box = slide.shapes.add_textbox(x, y, width, Inches(2.2))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        if items:
            for i, item in enumerate(items[:max_items]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = f"• {item}"
                p.font.size = Pt(14)
                p.space_after = Pt(8)
                
                if highlight_all:
                    p.font.color.rgb = MAGENTA
                    p.font.bold = True
                else:
                    p.font.color.rgb = DARK_TEXT
        else:
            tf.paragraphs[0].text = "• (To be discovered)"
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
            tf.paragraphs[0].font.italic = True
    
    # ===== KEY INITIATIVES (Left Column) =====
    add_section_title(slide1, left_x, top_y, "Key Initiatives:")
    add_bullet_points(slide1, left_x, top_y + Inches(0.5), data.key_initiatives)
    
    # ===== CHALLENGES (Right Column) =====
    add_section_title(slide1, right_x, top_y, "Challenges:")
    add_bullet_points(slide1, right_x, top_y + Inches(0.5), data.challenges)
    
    # ===== SOLUTION REQUIREMENTS (Full Width) =====
    req_y = Inches(4.2)
    add_section_title(slide1, left_x, req_y, "Solution Requirements:", width=Inches(12.333))
    
    # Solution requirements as bullet points
    req_box = slide1.shapes.add_textbox(left_x, req_y + Inches(0.5), Inches(12.333), Inches(2.2))
    tf = req_box.text_frame
    tf.word_wrap = True
    
    if data.solution_requirements:
        for i, item in enumerate(data.solution_requirements[:6]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(6)
    else:
        tf.paragraphs[0].text = "• (To be discovered)"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        tf.paragraphs[0].font.italic = True
    
    _add_footer(slide1)
    
    # =========================================
    # SLIDE 2: QUESTIONS FOR NEXT CALL
    # =========================================
    slide2 = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide2)
    
    # Title
    q_title_text = "🎯 Questions for Next Call"
    if data.customer_name:
        q_title_text = f"🎯 Questions for Next Call with {data.customer_name}"
    
    q_title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tf = q_title_box.text_frame
    tf.paragraphs[0].text = q_title_text
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MAGENTA
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.4))
    tf = subtitle_box.text_frame
    tf.paragraphs[0].text = "Probing questions to uncover gaps and deepen discovery"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = DARK_TEXT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Questions - large and prominent
    questions_box = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.5))
    tf = questions_box.text_frame
    tf.word_wrap = True
    
    if data.follow_up_questions:
        for i, question in enumerate(data.follow_up_questions[:6]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # Number the questions
            p.text = f"{i + 1}. {question}"
            p.font.size = Pt(18)
            p.font.color.rgb = MAGENTA
            p.font.bold = True
            p.space_after = Pt(20)
    else:
        tf.paragraphs[0].text = "No specific questions identified - review missed opportunities in the analysis."
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        tf.paragraphs[0].font.italic = True
    
    _add_footer(slide2)
    
    # Save to bytes
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    return pptx_bytes.getvalue()


def generate_recap_slide(data: RecapSlideData) -> bytes:
    """
    Convenience function to generate a recap presentation.
    
    Args:
        data: RecapSlideData containing the content
        
    Returns:
        Bytes of the PowerPoint file
    """
    return create_recap_presentation(data)

